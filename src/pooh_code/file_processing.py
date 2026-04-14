"""文件处理管线：将用户上传的文件转换为 LLM 可消费的格式。

- 图片 → base64 data URL（发送给多模态 LLM）
- PDF → 提取文本
- Office 文档 → 提取文本
- 视频 → 提取关键帧为图片
- 其他 → 尝试读取为纯文本
"""

from __future__ import annotations

import base64
import logging
import mimetypes
import subprocess
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# 图片 MIME 类型
IMAGE_MIMES = {"image/png", "image/jpeg", "image/gif", "image/webp", "image/svg+xml"}
VIDEO_MIMES = {"video/mp4", "video/quicktime", "video/x-msvideo", "video/webm", "video/mpeg"}

# 文件后缀 → 处理策略
_STRATEGY = {
    ".png": "image", ".jpg": "image", ".jpeg": "image",
    ".gif": "image", ".webp": "image", ".bmp": "image",
    ".pdf": "pdf",
    ".docx": "docx", ".doc": "docx",
    ".xlsx": "xlsx", ".xls": "xlsx", ".csv": "csv", ".tsv": "csv",
    ".pptx": "pptx", ".ppt": "pptx",
    ".mp4": "video", ".mov": "video", ".avi": "video", ".webm": "video",
    ".txt": "text", ".md": "text", ".json": "text", ".py": "text",
    ".js": "text", ".ts": "text", ".html": "text", ".css": "text",
    ".xml": "text", ".yaml": "text", ".yml": "text", ".toml": "text",
    ".sh": "text", ".sql": "text", ".log": "text",
}


def detect_strategy(file_path: Path) -> str:
    return _STRATEGY.get(file_path.suffix.lower(), "binary")


def process_file(file_path: Path) -> list[dict[str, Any]]:
    """处理单个文件，返回 LLM 消息内容块列表。

    返回格式：
    - 图片: [{"type": "image", "media_type": "image/png", "data": "<base64>"}]
    - 文本/文档: [{"type": "text", "text": "<提取的内容>"}]
    - 视频: [{"type": "image", ...}, {"type": "text", "text": "视频关键帧"}] （多帧）
    """
    strategy = detect_strategy(file_path)

    if strategy == "image":
        return _process_image(file_path)
    elif strategy == "pdf":
        return _process_pdf(file_path)
    elif strategy == "docx":
        return _process_docx(file_path)
    elif strategy == "xlsx":
        return _process_xlsx(file_path)
    elif strategy == "csv":
        return _process_csv(file_path)
    elif strategy == "pptx":
        return _process_pptx(file_path)
    elif strategy == "video":
        return _process_video(file_path)
    elif strategy == "text":
        return _process_text(file_path)
    else:
        return [{"type": "text", "text": f"[已上传文件: {file_path.name}（{file_path.suffix} 格式，无法直接解析内容）]"}]


def _process_image(file_path: Path) -> list[dict[str, Any]]:
    """图片 → base64 编码，发给多模态 LLM。"""
    data = file_path.read_bytes()
    b64 = base64.standard_b64encode(data).decode("ascii")
    mime = mimetypes.guess_type(str(file_path))[0] or "image/png"
    return [
        {
            "type": "text",
            "text": f"🖼️ 图片文件: {file_path.name} ({_human_size(file_path.stat().st_size)})",
        },
        {
            "type": "image",
            "media_type": mime,
            "data": b64,
            "filename": file_path.name,
        },
    ]


def _process_pdf(file_path: Path) -> list[dict[str, Any]]:
    """PDF → 提取文本内容。"""
    text = ""
    # 尝试 PyPDF2
    try:
        import PyPDF2
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            pages = []
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    pages.append(f"--- 第 {i + 1} 页 ---\n{page_text.strip()}")
            text = "\n\n".join(pages)
    except ImportError:
        pass

    # 尝试 pdfplumber（更好的表格支持）
    if not text:
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                pages = []
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        pages.append(f"--- 第 {i + 1} 页 ---\n{page_text.strip()}")
                text = "\n\n".join(pages)
        except ImportError:
            pass

    if not text:
        return [{"type": "text", "text": f"[已上传 PDF: {file_path.name}，但未安装 PDF 解析库（PyPDF2 或 pdfplumber）。请运行 `uv add PyPDF2` 安装。]"}]

    return [{"type": "text", "text": f"📄 PDF 文件: {file_path.name}\n\n{_truncate(text)}"}]


def _process_docx(file_path: Path) -> list[dict[str, Any]]:
    """Word 文档 → 提取文本。"""
    try:
        from docx import Document
        doc = Document(str(file_path))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        # 也提取表格内容
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                parts.append("\n".join(rows))
        text = "\n\n".join(parts)
    except ImportError:
        return [{"type": "text", "text": f"[已上传 Word 文档: {file_path.name}，但未安装 python-docx。请运行 `uv add python-docx` 安装。]"}]
    except Exception as exc:
        return [{"type": "text", "text": f"[解析 Word 文档失败: {exc}]"}]

    return [{"type": "text", "text": f"📄 Word 文档: {file_path.name}\n\n{_truncate(text)}"}]


def _process_xlsx(file_path: Path) -> list[dict[str, Any]]:
    """Excel 文件 → 提取表格数据。"""
    try:
        import pandas as pd
        sheets = pd.read_excel(file_path, sheet_name=None, dtype=str)
        parts = []
        for sheet_name, df in sheets.items():
            header = f"### Sheet: {sheet_name} ({len(df)} 行 × {len(df.columns)} 列)"
            table_str = df.head(100).to_string(index=False)
            parts.append(f"{header}\n{table_str}")
            if len(df) > 100:
                parts.append(f"... 共 {len(df)} 行，仅显示前 100 行")
        text = "\n\n".join(parts)
    except ImportError:
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True, data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(max_row=101, values_only=True):
                    rows.append(" | ".join(str(c) if c is not None else "" for c in row))
                parts.append(f"### Sheet: {sheet_name}\n" + "\n".join(rows))
            wb.close()
            text = "\n\n".join(parts)
        except ImportError:
            return [{"type": "text", "text": f"[已上传 Excel: {file_path.name}，但未安装 pandas 或 openpyxl。请运行 `uv add pandas openpyxl` 安装。]"}]
    except Exception as exc:
        return [{"type": "text", "text": f"[解析 Excel 文件失败: {exc}]"}]

    return [{"type": "text", "text": f"📊 Excel 文件: {file_path.name}\n\n{_truncate(text)}"}]


def _process_csv(file_path: Path) -> list[dict[str, Any]]:
    """CSV/TSV → 读取为文本。"""
    try:
        import pandas as pd
        sep = "\t" if file_path.suffix.lower() == ".tsv" else ","
        df = pd.read_csv(file_path, sep=sep, dtype=str, nrows=200)
        text = f"表格: {len(df)} 行 × {len(df.columns)} 列\n\n{df.head(100).to_string(index=False)}"
    except ImportError:
        text = file_path.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        return [{"type": "text", "text": f"[解析 CSV 失败: {exc}]"}]

    return [{"type": "text", "text": f"📊 {file_path.name}\n\n{_truncate(text)}"}]


def _process_pptx(file_path: Path) -> list[dict[str, Any]]:
    """PowerPoint → 提取文本。"""
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        texts.append(" | ".join(cells))
            if texts:
                parts.append(f"--- 幻灯片 {i} ---\n" + "\n".join(texts))
        text = "\n\n".join(parts)
    except ImportError:
        return [{"type": "text", "text": f"[已上传 PPT: {file_path.name}，但未安装 python-pptx。请运行 `uv add python-pptx` 安装。]"}]
    except Exception as exc:
        return [{"type": "text", "text": f"[解析 PPT 失败: {exc}]"}]

    return [{"type": "text", "text": f"📑 PPT 文件: {file_path.name}\n\n{_truncate(text)}"}]


def _process_video(file_path: Path) -> list[dict[str, Any]]:
    """视频 → 提取关键帧 + 语音转文字。"""
    results: list[dict[str, Any]] = []
    results.append({
        "type": "text",
        "text": f"🎬 视频文件: {file_path.name} ({_human_size(file_path.stat().st_size)})",
    })

    tmp_dir = file_path.parent / f".{file_path.stem}_frames"

    # 1) 提取关键帧
    try:
        tmp_dir.mkdir(exist_ok=True)
        subprocess.run(
            [
                "ffmpeg", "-i", str(file_path),
                "-vf", "fps=1/10,scale=640:-1",  # 每 10 秒一帧，缩放到 640px 宽
                "-frames:v", "4",
                "-q:v", "2",
                str(tmp_dir / "frame_%02d.jpg"),
            ],
            capture_output=True, timeout=30,
        )
        frames = sorted(tmp_dir.glob("frame_*.jpg"))
        for frame in frames[:4]:
            results.extend(_process_image(frame))
        for frame in frames:
            frame.unlink(missing_ok=True)
    except FileNotFoundError:
        results.append({"type": "text", "text": "[未安装 ffmpeg，无法提取视频帧和音频。]"})
        return results
    except Exception as exc:
        results.append({"type": "text", "text": f"[提取视频帧失败: {exc}]"})

    # 2) 提取音频并语音转文字
    try:
        audio_path = tmp_dir / "audio.wav"
        tmp_dir.mkdir(exist_ok=True)
        proc = subprocess.run(
            [
                "ffmpeg", "-i", str(file_path),
                "-vn",                    # 不要视频
                "-acodec", "pcm_s16le",   # 16-bit PCM
                "-ar", "16000",           # 16kHz 采样率（whisper 要求）
                "-ac", "1",               # 单声道
                "-y",                     # 覆盖
                str(audio_path),
            ],
            capture_output=True, timeout=60,
        )
        if audio_path.exists() and audio_path.stat().st_size > 1000:
            transcript = _transcribe_audio(audio_path)
            if transcript.strip():
                results.append({
                    "type": "text",
                    "text": f"🗣️ 视频语音转文字:\n\n{_truncate(transcript, 15000)}",
                })
            else:
                results.append({"type": "text", "text": "[视频中未检测到语音内容]"})
            audio_path.unlink(missing_ok=True)
        else:
            results.append({"type": "text", "text": "[视频中无音频轨道]"})
    except Exception as exc:
        results.append({"type": "text", "text": f"[提取视频音频失败: {exc}]"})

    # 清理临时目录
    try:
        if tmp_dir.exists():
            for f in tmp_dir.iterdir():
                f.unlink(missing_ok=True)
            tmp_dir.rmdir()
    except Exception:
        pass

    return results


def _transcribe_audio(audio_path: Path) -> str:
    """使用 faster-whisper 本地模型将音频转为文字。"""
    try:
        from faster_whisper import WhisperModel
    except ImportError:
        return "[未安装 faster-whisper，无法转录语音。请运行 `uv add faster-whisper` 安装。]"

    model = WhisperModel("base", device="cpu", compute_type="int8")
    segments, info = model.transcribe(str(audio_path), beam_size=5)
    logger.info("检测到语言: %s (概率 %.2f)", info.language, info.language_probability)

    lines: list[str] = []
    for seg in segments:
        ts = f"[{_fmt_ts(seg.start)} -> {_fmt_ts(seg.end)}]"
        lines.append(f"{ts} {seg.text.strip()}")
    return "\n".join(lines)


def _fmt_ts(seconds: float) -> str:
    """格式化时间戳为 mm:ss。"""
    m, s = divmod(int(seconds), 60)
    return f"{m:02d}:{s:02d}"


def _process_text(file_path: Path) -> list[dict[str, Any]]:
    """纯文本文件 → 直接读取。"""
    try:
        text = file_path.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        return [{"type": "text", "text": f"[读取文件失败: {exc}]"}]
    return [{"type": "text", "text": f"📝 {file_path.name}\n\n{_truncate(text)}"}]


def _truncate(text: str, limit: int = 30000) -> str:
    if len(text) <= limit:
        return text
    return text[:limit] + f"\n\n[... 内容过长，已截断，共 {len(text)} 字符]"


def _human_size(size: int) -> str:
    if size < 1024:
        return f"{size} B"
    if size < 1024 * 1024:
        return f"{size / 1024:.1f} KB"
    return f"{size / (1024 * 1024):.1f} MB"
